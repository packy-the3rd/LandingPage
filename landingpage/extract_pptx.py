"""
PPTXファイルからテキストと画像を抽出するスクリプト
"""
import os
import sys
import io

# Windows cp932 エンコーディング問題の対策
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE


def extract_pptx(pptx_path, output_dir=None):
    """PPTXファイルからテキストと画像を抽出する"""

    if not os.path.exists(pptx_path):
        print(f"エラー: ファイルが見つかりません: {pptx_path}")
        return

    # 出力ディレクトリ
    if output_dir is None:
        output_dir = os.path.join(os.path.dirname(pptx_path), "extracted")
    os.makedirs(output_dir, exist_ok=True)

    images_dir = os.path.join(output_dir, "images")
    os.makedirs(images_dir, exist_ok=True)

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"エラー: PPTXファイルを開けませんでした: {e}")
        return

    print(f"ファイル: {pptx_path}")
    print(f"スライド数: {len(prs.slides)}")
    print(f"スライドサイズ: {prs.slide_width} x {prs.slide_height}")
    print("=" * 60)

    all_text = []
    image_count = 0

    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_text = []
        print(f"\n--- スライド {slide_idx} ---")

        for shape in slide.shapes:
            # テキストの抽出
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        print(f"  テキスト: {text}")
                        slide_text.append(text)

            # テーブルの抽出
            if shape.has_table:
                table = shape.table
                print(f"  テーブル ({table.rows.__len__()} 行 x {len(table.columns)} 列):")
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    print(f"    | {' | '.join(row_text)} |")
                    slide_text.append(" | ".join(row_text))

            # 画像の抽出
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                ext = image.content_type.split("/")[-1]
                if ext == "jpeg":
                    ext = "jpg"
                image_filename = f"slide{slide_idx}_img{image_count + 1}.{ext}"
                image_path = os.path.join(images_dir, image_filename)
                with open(image_path, "wb") as f:
                    f.write(image.blob)
                image_count += 1
                print(f"  画像保存: {image_filename}")

            # グループ内の画像も抽出
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for child_shape in shape.shapes:
                    if child_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image = child_shape.image
                        ext = image.content_type.split("/")[-1]
                        if ext == "jpeg":
                            ext = "jpg"
                        image_filename = f"slide{slide_idx}_img{image_count + 1}.{ext}"
                        image_path = os.path.join(images_dir, image_filename)
                        with open(image_path, "wb") as f:
                            f.write(image.blob)
                        image_count += 1
                        print(f"  画像保存 (グループ内): {image_filename}")

                    if child_shape.has_text_frame:
                        for paragraph in child_shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                print(f"  テキスト (グループ内): {text}")
                                slide_text.append(text)

        all_text.append({"slide": slide_idx, "texts": slide_text})

    # テキストをファイルに保存
    text_file = os.path.join(output_dir, "extracted_text.txt")
    with open(text_file, "w", encoding="utf-8") as f:
        for slide_data in all_text:
            f.write(f"=== スライド {slide_data['slide']} ===\n")
            for text in slide_data["texts"]:
                f.write(f"{text}\n")
            f.write("\n")

    print("\n" + "=" * 60)
    print(f"抽出完了!")
    print(f"  テキスト: {text_file}")
    print(f"  画像数: {image_count} 枚 → {images_dir}")


if __name__ == "__main__":
    pptx_path = r"C:\Users\ysweb\Dropbox\12.広報\LP企画\LP企画v1.pptx"
    extract_pptx(pptx_path)
